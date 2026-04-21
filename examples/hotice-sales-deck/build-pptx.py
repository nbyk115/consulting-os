"""
Build hotice-deck.pptx from the rendered slide PNG files.
Each slide is a full-bleed 1280x720 image → zero font/layout risk.
"""
import os
from pptx import Presentation
from pptx.util import Emu, Inches

HERE = os.path.dirname(os.path.abspath(__file__))

# Slide order matches the deck narrative (cover → P2-P18, including s-prefixed legacy ids)
ORDER = [
    ("cover", "P1 — Cover"),
    ("p2",    "P2 — Who is hotice"),
    ("p3",    "P3 — Market"),
    ("s2",    "P4 — Influencer Marketing 3.0"),
    ("p5",    "P5 — Agency Pains"),
    ("p7",    "P6 — Translator Positioning"),
    ("s3",    "P7 — Competitive Landscape"),
    ("p8a",   "P8 — Why hotice"),
    ("p9",    "P9 — Global Implementation"),
    ("p10",   "P10 — Speed x Quality"),
    ("s1",    "P11 — Proof Summary"),
    ("p12",   "P12 — Case Shiseido"),
    ("p13",   "P13 — Case Amazon EC"),
    ("p14",   "P14 — Case New Cafe"),
    ("p15",   "P15 — Case D2C"),
    ("p16",   "P16 — Partnership Modes"),
    ("p17",   "P17 — Pricing / Team / Delivery"),
    ("p18",   "P18 — CTA"),
]

prs = Presentation()
# 16:9 slide size, match our 1280x720 canvas aspect ratio
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # Blank layout

for sid, label in ORDER:
    png_path = os.path.join(HERE, f"slide-{sid}.png")
    if not os.path.exists(png_path):
        print(f"! missing {png_path}")
        continue
    slide = prs.slides.add_slide(blank_layout)
    # Full-bleed image
    slide.shapes.add_picture(png_path, 0, 0,
                             width=prs.slide_width,
                             height=prs.slide_height)
    # Alt text for accessibility
    pic = slide.shapes[-1]
    pic.name = label
    print(f"+ {label}")

out_path = os.path.join(HERE, "hotice-deck.pptx")
prs.save(out_path)
size_kb = os.path.getsize(out_path) // 1024
print(f"\n✅ Saved: {out_path} ({size_kb} KB, {len(prs.slides)} slides)")
